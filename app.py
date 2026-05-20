import streamlit as st


import pandas as pd
import matplotlib.pyplot as plt
import folium
import plotly.express as px
import sqlite3
import pdfplumber
import re
import json
import os
import shutil

from datetime import datetime

from datetime import datetime
from geopy.geocoders import Nominatim
from io import BytesIO
from streamlit_folium import st_folium

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from reportlab.platypus import (
    SimpleDocTemplate,
    Paragraph,
    Spacer,
    Table,
    TableStyle,
    Image
)
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4

from dotenv import load_dotenv


load_dotenv()

# =========================
# CONFIGURAÇÃO DO BANCO
# =========================

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DB_NAME = "crm_pedidos.db"

def conectar():
    return sqlite3.connect(DB_NAME)

def criar_backup():
    data = datetime.now().strftime("%Y%m%d_%H%M")
    shutil.copy(
        "crm_pedidos.db",
        f"backup/backup_{data}.db"
    )


def criar_tabelas():
    conn = conectar()
    cursor = conn.cursor()

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS clientes (
            codigo_cliente TEXT PRIMARY KEY,
            nome TEXT,
            email TEXT,
            telefone TEXT,
            cnpj TEXT,
            endereco TEXT,
            cidade TEXT,
            estado TEXT,
            latitude REAL,
            longitude REAL,
            metodo_pagamento TEXT
        )
    """)

    try:
        cursor.execute("ALTER TABLE clientes ADD COLUMN latitude REAL")
    except sqlite3.OperationalError:
        pass

    try:
        cursor.execute("ALTER TABLE clientes ADD COLUMN longitude REAL")
    except sqlite3.OperationalError:
        pass

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS pedidos (
            numero_pedido TEXT PRIMARY KEY,
            codigo_cliente TEXT,
            valor_total REAL,
            origem TEXT,
            FOREIGN KEY (codigo_cliente) REFERENCES clientes(codigo_cliente)
        )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS itens_pedido (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            numero_pedido TEXT,
            codigo_produto TEXT,
            produto TEXT,
            quantidade REAL,
            preco_unitario REAL,
            total_sem_st REAL,
            st REAL,
            total_com_st REAL,
            FOREIGN KEY (numero_pedido) REFERENCES pedidos(numero_pedido)
        )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS pipeline (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            codigo_cliente TEXT,
            cliente TEXT,
            etapa TEXT,
            responsavel TEXT,
            observacao TEXT,
            data_atualizacao TEXT
        )
    """)

    cursor.execute("""
        CREATE TABLE IF NOT EXISTS estoque_diario (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            data_importacao TEXT,
            arquivo_origem TEXT,
            aba_origem TEXT,
            centro TEXT,
            material TEXT,
            descricao TEXT,
            categoria TEXT,
            estoque_total REAL,
            dados_json TEXT
        )
    """)

    conn.commit()
    conn.close()

# =========================
# LEITOR DE PDF DE PEDIDO
# =========================

def limpar_valor(valor):
    if not valor:
        return 0.0
    valor = valor.replace("R$", "").replace(".", "").replace(",", ".").strip()
    try:
        return float(valor)
    except:
        return 0.0


def extrair_texto_pdf(arquivo_pdf):
    texto = ""
    with pdfplumber.open(arquivo_pdf) as pdf:
        for pagina in pdf.pages:
            texto_extraido = pagina.extract_text()
            if texto_extraido:
                texto += texto_extraido + "\n"
    return texto


def geocodificar_endereco(endereco):
    if not endereco:
        return None, None

    try:
        geolocator = Nominatim(user_agent="ea_comercial_intelligence")
        location = geolocator.geocode(endereco, timeout=10)

        if location:
            return location.latitude, location.longitude

    except Exception:
        pass

    return None, None


def identificar_coluna(df_base, opcoes):
    for coluna in opcoes:
        if coluna in df_base.columns:
            return coluna
    return None


def preparar_estoque(df_base):
    df_base = df_base.copy()
    df_base.columns = df_base.columns.astype(str).str.strip()
    df_base = df_base.dropna(how="all")

    coluna_material = identificar_coluna(
        df_base,
        ["Material", "Codigo", "Código", "Cod Material", "Código Material"]
    )

    coluna_descricao = identificar_coluna(
        df_base,
        ["Texto breve de material", "Produto", "Descrição", "Descricao", "Material Descrição"]
    )

    coluna_categoria = identificar_coluna(
        df_base,
        ["Categorias", "Categoria", "Grupo", "Linha"]
    )

    coluna_centro = identificar_coluna(
        df_base,
        ["Centro", "CD", "Unidade"]
    )

    coluna_total = identificar_coluna(
        df_base,
        ["Total", "Estoque", "Estoque Total", "Saldo", "Qtd", "Quantidade"]
    )

    if coluna_material is None:
        df_base["Material"] = ""
        coluna_material = "Material"

    if coluna_descricao is None:
        df_base["Texto breve de material"] = ""
        coluna_descricao = "Texto breve de material"

    if coluna_categoria is None:
        df_base["Categorias"] = "Não informado"
        coluna_categoria = "Categorias"

    if coluna_centro is None:
        df_base["Centro"] = "Não informado"
        coluna_centro = "Centro"

    if coluna_total is None:
        possiveis_numericas = []
        for coluna in df_base.columns:
            serie = pd.to_numeric(df_base[coluna], errors="coerce")
            if serie.notna().sum() > 0:
                possiveis_numericas.append(coluna)

        if possiveis_numericas:
            coluna_total = possiveis_numericas[-1]
        else:
            df_base["Total"] = 0
            coluna_total = "Total"

    df_base["_material"] = df_base[coluna_material].astype(str)
    df_base["_descricao"] = df_base[coluna_descricao].astype(str)
    df_base["_categoria"] = df_base[coluna_categoria].astype(str)
    df_base["_centro"] = df_base[coluna_centro].astype(str)
    df_base["_estoque_total"] = pd.to_numeric(
        df_base[coluna_total],
        errors="coerce"
    ).fillna(0)

    return df_base, {
        "material": coluna_material,
        "descricao": coluna_descricao,
        "categoria": coluna_categoria,
        "centro": coluna_centro,
        "total": coluna_total,
    }


def salvar_estoque_no_banco(df_base, nome_arquivo, aba_origem):
    conn = conectar()
    cursor = conn.cursor()

    cursor.execute("DELETE FROM estoque_diario")

    data_importacao = datetime.now().strftime("%d/%m/%Y %H:%M")

    for _, row in df_base.iterrows():
        cursor.execute("""
            INSERT INTO estoque_diario (
                data_importacao,
                arquivo_origem,
                aba_origem,
                centro,
                material,
                descricao,
                categoria,
                estoque_total,
                dados_json
            ) VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """, (
            data_importacao,
            nome_arquivo,
            aba_origem,
            row.get("_centro", ""),
            row.get("_material", ""),
            row.get("_descricao", ""),
            row.get("_categoria", ""),
            float(row.get("_estoque_total", 0)),
            json.dumps(row.to_dict(), ensure_ascii=False, default=str)
        ))

    conn.commit()
    conn.close()

    return data_importacao

def limpar_texto(valor):
    if not valor:
        return ""
    return " ".join(str(valor).split()).strip()


def extrair_itens_pedido(texto):
    itens = []
    linhas = [linha.strip() for linha in texto.split("\n") if linha.strip()]

    item_atual = None

    for linha in linhas:

        # Linha de valores: Quant. Preço Total ST Total c/ ST
        valores = re.findall(r"R\$\s*([\d.,]+)", linha)

        if valores and item_atual and len(valores) >= 4:

            preco_unitario = limpar_valor(valores[0])
            total_sem_st = limpar_valor(valores[1])
            st = limpar_valor(valores[2])
            total_com_st = limpar_valor(valores[3])

            # Calcula quantidade real pelo total / preço para evitar erro quando o PDF quebra a quantidade em outra linha
            if preco_unitario > 0:
                quantidade = round(total_sem_st / preco_unitario, 2)
            else:
                quantidade = 0

            itens.append({
                "codigo_produto": item_atual["codigo_produto"],
                "produto": limpar_texto(item_atual["produto"]),
                "quantidade": quantidade,
                "preco_unitario": preco_unitario,
                "total_sem_st": total_sem_st,
                "st": st,
                "total_com_st": total_com_st,
            })

            item_atual = None
            continue

        # Linha começando com código do produto
        match_codigo = re.match(r"^(\d{5,6})\s*(.*)", linha)

        if match_codigo:
            item_atual = {
                "codigo_produto": match_codigo.group(1),
                "produto": match_codigo.group(2)
            }
            continue

        # Continuação da descrição do produto
        if item_atual:
            if not linha.startswith("Codigo Produto") and not linha.startswith("Quant."):
                linha_limpa = linha

                # Corrige caso venha algo tipo "5 12KG"
                if re.match(r"^\d+\s+\d+(KG|G|UN|CX)", linha_limpa, re.IGNORECASE):
                    partes = linha_limpa.split(" ", 1)
                    linha_limpa = partes[1]

                item_atual["produto"] += " " + linha_limpa

    return itens


def extrair_dados_pedido(arquivo_pdf):
    texto = extrair_texto_pdf(arquivo_pdf)

    pedido = re.search(r"Pedido SAP\s+(\d+)", texto)
    codigo_cliente = re.search(r"Codigo do Cliente\s+(\d+)", texto)
    nome = re.search(r"Nome\s+(.+)", texto)
    email = re.search(r"E-mail\s+(.+)", texto)
    cnpj = re.search(r"CNPJ\s+(\d+)", texto)

    endereco = re.search(
        r"Endere[cç]o de Entrega\s+(.+?)\s+Metodo de Pagamento",
        texto,
        re.DOTALL
    )

    metodo_pagamento = re.search(
        r"Metodo de Pagamento\s+(.+?)\s+Itens do Pedido",
        texto,
        re.DOTALL
    )

    bloco_resumo = re.search(
        r"Resumo de Totais(.*)",
        texto,
        re.DOTALL
    )

    valor_total = 0.0

    if bloco_resumo:
        valores = re.findall(
            r"R\$\s*([\d.,]+)",
            bloco_resumo.group(1)
        )

        if valores:
            valor_total = limpar_valor(valores[-1])

    itens = extrair_itens_pedido(texto)

    return {
        "numero_pedido": pedido.group(1) if pedido else "",
        "codigo_cliente": codigo_cliente.group(1) if codigo_cliente else "",
        "nome": limpar_texto(nome.group(1)) if nome else "",
        "email": limpar_texto(email.group(1)) if email else "",
        "cnpj": limpar_texto(cnpj.group(1)) if cnpj else "",
        "endereco": limpar_texto(endereco.group(1)) if endereco else "",
        "metodo_pagamento": limpar_texto(metodo_pagamento.group(1)) if metodo_pagamento else "",
        "valor_total": valor_total,
        "itens": itens,
    }

# =========================
# CONFIG STREAMLIT
# =========================

st.set_page_config(
    page_title="EA Comercial Intelligence",
    page_icon="📊",
    layout="wide"
)

criar_tabelas()

# =========================
# CSS PREMIUM — EA COMERCIAL INTELLIGENCE
# =========================

st.markdown("""
<style>

/* =========================
BASE / FUNDO
========================= */

#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}

html, body, .stApp {
    background: #0f172a !important;
    color: #f8fafc !important;
    font-family: 'Segoe UI', sans-serif !important;
}

[data-testid="stAppViewContainer"] {
    background:
        radial-gradient(circle at top left, rgba(179,0,0,0.12), transparent 30%),
        linear-gradient(180deg, #0f172a 0%, #111827 100%) !important;
}

.main .block-container,
.block-container,
[data-testid="stMainBlockContainer"] {
    background: transparent !important;
    color: #f8fafc !important;
    max-width: 1450px !important;
    padding-top: 1.2rem !important;
    padding-bottom: 2rem !important;
}

/* =========================
TEXTOS
========================= */

h1, h2, h3, h4, h5, h6,
p, label, span, small,
[data-testid="stMarkdownContainer"] {
    color: #f8fafc !important;
}

.ea-section-title {
    font-size: 28px !important;
    font-weight: 800 !important;
    color: #ffffff !important;
    margin-top: 20px !important;
    margin-bottom: 18px !important;
    letter-spacing: -0.5px !important;
}


/* =========================
SIDEBAR
========================= */

[data-testid="stSidebar"] {
    background:
        linear-gradient(180deg, #020617 0%, #071226 100%) !important;
    border-right: 1px solid rgba(255,255,255,0.08) !important;
}

[data-testid="stSidebar"] * {
    color: #f8fafc !important;
}

[data-testid="stSidebar"] img {
    margin-top: 12px !important;
    margin-bottom: 10px !important;
}

[data-testid="stSidebar"] .stMultiSelect,
[data-testid="stSidebar"] [data-baseweb="select"] {
    background: rgba(15,23,42,0.92) !important;
    border-radius: 10px !important;
}

/* =========================
ABAS
========================= */

.stTabs [data-baseweb="tab-list"] {
    gap: 8px !important;
    border-bottom: 1px solid rgba(255,255,255,0.08) !important;
}

.stTabs [data-baseweb="tab"] {
    background: #1f2937 !important;
    color: #d1d5db !important;
    border-radius: 12px 12px 0 0 !important;
    padding-left: 18px !important;
    padding-right: 18px !important;
    height: 46px !important;
    font-weight: 700 !important;
    border: 1px solid rgba(255,255,255,0.05) !important;
}

.stTabs [aria-selected="true"],
button[data-baseweb="tab"][aria-selected="true"] {
    background: #b30000 !important;
    color: white !important;
    border-color: rgba(255,255,255,0.15) !important;
}

/* =========================
CARDS / MÉTRICAS PREMIUM
========================= */

[data-testid="metric-container"],
div[data-testid="stMetric"] {
    background:
        linear-gradient(145deg, rgba(17,24,39,0.98), rgba(30,41,59,0.98)) !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    border-radius: 18px !important;
    padding: 18px !important;
    box-shadow: 0 8px 26px rgba(0,0,0,0.30) !important;
    transition: all 0.25s ease !important;
}

[data-testid="metric-container"]:hover,
div[data-testid="stMetric"]:hover {
    transform: translateY(-2px) !important;
    border-color: rgba(179,0,0,0.45) !important;
}

[data-testid="stMetricLabel"],
[data-testid="metric-container"] label {
    color: #9ca3af !important;
    font-size: 13px !important;
    font-weight: 600 !important;
}

[data-testid="stMetricValue"],
[data-testid="metric-container"] [data-testid="stMetricValue"] {
    color: #ffffff !important;
    font-size: 28px !important;
    font-weight: 800 !important;
}

[data-testid="metric-container"] p {
    color: #d1d5db !important;
}

/* =========================
ALERTAS
========================= */

.stAlert {
    border-radius: 14px !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    box-shadow: 0 4px 18px rgba(0,0,0,0.18) !important;
}

/* =========================
FORMULÁRIOS / INPUTS
========================= */

textarea,
div[data-baseweb="select"] > div {
    background: #111827 !important;
    color: #f8fafc !important;
    border: 1px solid rgba(255,255,255,0.16) !important;
    border-radius: 10px !important;
}

textarea {
    min-height: 110px !important;
}

div[data-baseweb="select"] * {
    color: #f8fafc !important;
}

.stButton > button,
button[kind="primaryFormSubmit"],
button[kind="secondaryFormSubmit"] {
    background: #b30000 !important;
    color: #ffffff !important;
    border-radius: 10px !important;
    border: 1px solid rgba(255,255,255,0.10) !important;
    font-weight: 700 !important;
    padding: 0.45rem 1rem !important;
    box-shadow: 0 4px 14px rgba(179,0,0,0.28) !important;
}

.stButton > button:hover,
button[kind="primaryFormSubmit"]:hover,
button[kind="secondaryFormSubmit"]:hover {
    background: #d00000 !important;
    color: white !important;
    border-color: rgba(255,255,255,0.25) !important;
}

/* =========================
DATAFRAMES
========================= */

[data-testid="stDataFrame"] {
    background: #ffffff !important;
    border-radius: 14px !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    overflow: hidden !important;
    box-shadow: 0 6px 22px rgba(0,0,0,0.22) !important;
}

[data-testid="stDataFrame"] * {
    color: #111827 !important;
}

/* =========================
PLOTLY / GRÁFICOS
========================= */

.js-plotly-plot {
    border-radius: 16px !important;
    overflow: hidden !important;
}

.js-plotly-plot .plotly .main-svg {
    background: transparent !important;
}

/* =========================
CARDS HTML CUSTOMIZADOS
========================= */

.ea-card {
    background: rgba(15,23,42,0.88) !important;
    border-radius: 18px !important;
    padding: 22px !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    margin-bottom: 18px !important;
    box-shadow: 0 6px 22px rgba(0,0,0,0.22) !important;
}

/* =========================
MOBILE
========================= */

@media(max-width:900px) {
    .block-container,
    [data-testid="stMainBlockContainer"] {
        padding-left: 1rem !important;
        padding-right: 1rem !important;
    }

    .stTabs [data-baseweb="tab"] {
        padding-left: 10px !important;
        padding-right: 10px !important;
        font-size: 12px !important;
    }

    [data-testid="stMetricValue"] {
        font-size: 22px !important;
    }
}

input {
    color: white !important;
}

</style>
""", unsafe_allow_html=True)

# =========================
# LOGIN DO SISTEMA - MANUAL
# =========================

if "logado" not in st.session_state:
    st.session_state.logado = False

if not st.session_state.logado:
    st.markdown("## Login EA CRM")

    usuario = st.text_input("Usuário")
    senha = st.text_input("Senha", type="password")

    if st.button("Entrar"):
        if usuario == "euler" and senha == "123456":
            st.session_state.logado = True
            st.rerun()
        else:
            st.error("Usuário ou senha incorretos.")

    st.stop()

st.success("Bem-vindo, Euler Souza")

if st.button("Sair"):
    st.session_state.logado = False
    st.rerun()

# =========================
# TEMA GLOBAL DOS GRÁFICOS PLOTLY
# =========================

_original_plotly_chart = st.plotly_chart

def plotly_chart_premium(fig, *args, **kwargs):
    try:
        fig.update_layout(
            template="plotly_dark",
            paper_bgcolor="rgba(0,0,0,0)",
            plot_bgcolor="rgba(0,0,0,0)",
            font=dict(color="#f8fafc"),
            title_font=dict(color="#ffffff", size=18),
            legend=dict(
                bgcolor="rgba(0,0,0,0)",
                font=dict(color="#f8fafc")
            ),
            margin=dict(l=30, r=30, t=60, b=40)
        )
        fig.update_xaxes(
            gridcolor="rgba(255,255,255,0.08)",
            zerolinecolor="rgba(255,255,255,0.12)"
        )
        fig.update_yaxes(
            gridcolor="rgba(255,255,255,0.08)",
            zerolinecolor="rgba(255,255,255,0.12)"
        )
    except Exception:
        pass

    return _original_plotly_chart(fig, *args, **kwargs)

st.plotly_chart = plotly_chart_premium

# =========================
# DADOS GOOGLE SHEETS
# =========================

url = "https://docs.google.com/spreadsheets/d/1JQe8XfQfaPEOkHWHuQb7XvQ2aBuuC9UlwwBmZ-Lg2jI/export?format=csv"

df = pd.read_csv(url)

# limpar nomes das colunas
df.columns = df.columns.str.strip()

# =========================
# PADRONIZAR NOMES DA NOVA PLANILHA
# =========================

df.columns = df.columns.str.strip()

df = df.rename(columns={

    "Carimbo de data/hora": "Data",

    # CLIENTE
    "1. Nome do cliente": "Cliente",
    "1. Nome do cliente ": "Cliente",

    # RESPONSÁVEL
    "2. Nome do responsável pela compra": "Responsável Compra",
    "2. Nome do responsável pela compra ": "Responsável Compra",

    # WHATSAPP
    "3. WhatsApp para contato": "WhatsApp",
    "3. WhatsApp para contato ": "WhatsApp",

    # SEGMENTO
    "4. Segmento do cliente": "Segmento",
    "4. Segmento do cliente ": "Segmento",

    # CIDADE
    "5. Cidade": "Cidade",
    "5. Cidade ": "Cidade",

    # BAIRRO
    "6. Bairro": "Bairro",
    "6. Bairro ": "Bairro",

    # ENDEREÇO
    "7. Endereço": "Endereço",
    "7. Endereço ": "Endereço",

    # CEP
    "8. CEP": "CEP",
    "8. CEP ": "CEP",

    # FORNECEDOR
    "9. Quem é o principal fornecedor hoje?": "Concorrente",
    "9. Quem é o principal fornecedor hoje? ": "Concorrente",

    # TEMPO
    "10. Há quanto tempo trabalha com esse fornecedor?": "Tempo Fornecedor",
    "10. Há quanto tempo trabalha com esse fornecedor? ": "Tempo Fornecedor",

    # FATOR DECISÃO
    "11. O que mais pesa na decisão de compra?": "Fator Decisão",
    "11. O que mais pesa na decisão de compra? ": "Fator Decisão",

    # FORNECEDOR ATENDE
    "12. O fornecedor atual atende totalmente sua necessidade?": "Fornecedor Atende",
    "12. O fornecedor atual atende totalmente sua necessidade? ": "Fornecedor Atende",

    # DIFICULDADE FORNECEDOR
    "13. Qual principal dificuldade com o fornecedor atual?": "Dificuldade Fornecedor Atual",
    "13. Qual principal dificuldade com o fornecedor atual? ": "Dificuldade Fornecedor Atual",

    # PROBLEMAS
    "14. Quais problemas mais acontecem hoje?": "Problemas Atuais",
    "14. Quais problemas mais acontecem hoje? ": "Problemas Atuais",

    # QUALIDADE
    "15. Como você avalia a qualidade do fornecedor atual?": "Qualidade",
    "15. Como você avalia a qualidade do fornecedor atual? ": "Qualidade",

    # FREQUÊNCIA
    "16. Como você avalia a frequência das entregas?": "Frequência Entrega",
    "16. Como você avalia a frequência das entregas? ": "Frequência Entrega",

    # ATENDIMENTO
    "17. Como você avalia o atendimento comercial atual?": "Atendimento Atual",
    "17. Como você avalia o atendimento comercial atual? ": "Atendimento Atual",

    # JÁ COMPROU PIFPAF
    "18. Você já comprou da PifPaf?": "Ja Comprou PifPaf",
    "18. Você já comprou da PifPaf? ": "Ja Comprou PifPaf",

    # PIFPAF PREÇO
    "19. Como você avalia a PifPaf atualmente?  [Preço]": "PifPaf Preço",
    "19. Como você avalia a PifPaf atualmente? [Preço]": "PifPaf Preço",

    # PIFPAF QUALIDADE
    "19. Como você avalia a PifPaf atualmente?  [Qualidade]": "PifPaf Qualidade",
    "19. Como você avalia a PifPaf atualmente? [Qualidade]": "PifPaf Qualidade",

    # PIFPAF ENTREGA
    "19. Como você avalia a PifPaf atualmente?  [Entrega]": "PifPaf Entrega",
    "19. Como você avalia a PifPaf atualmente? [Entrega]": "PifPaf Entrega",

    # PIFPAF ATENDIMENTO
    "19. Como você avalia a PifPaf atualmente?  [Atendimento]": "PifPaf Atendimento",
    "19. Como você avalia a PifPaf atualmente? [Atendimento]": "PifPaf Atendimento",

    # PIFPAF VARIEDADE
    "19. Como você avalia a PifPaf atualmente?  [Variedade de produtos]": "PifPaf Variedade",
    "19. Como você avalia a PifPaf atualmente? [Variedade de produtos]": "PifPaf Variedade",

    # PIFPAF NEGOCIAÇÃO
    "19. Como você avalia a PifPaf atualmente?  [Negociação]": "PifPaf Negociação",
    "19. Como você avalia a PifPaf atualmente? [Negociação]": "PifPaf Negociação",

    # OPORTUNIDADE PIFPAF
    "20. O que faria você comprar mais da PifPaf?": "Oportunidade PifPaf",
    "20. O que faria você comprar mais da PifPaf? ": "Oportunidade PifPaf",

    # DIFICULDADE PIFPAF
    "21. Qual principal dificuldade com a PifPaf hoje?": "Dificuldade PifPaf",
    "21. Qual principal dificuldade com a PifPaf hoje? ": "Dificuldade PifPaf",

    # CATEGORIA
    "22. Qual categoria possui maior volume de vendas?": "Categoria Mais Vendida",
    "22. Qual categoria possui maior volume de vendas? ": "Categoria Mais Vendida",

    # VOLUME
    "23. Qual o volume médio mensal aproximado?": "Volume Mensal",
    "23. Qual o volume médio mensal aproximado? ": "Volume Mensal",

    # INTERESSE
    "Existe interesse em trocar fornecedor?": "Interesse em Trocar",

    # POTENCIAL
    "Potencial do cliente": "Potencial",

    # ESTADO
    "Estado": "Estado",

    # COORDENADAS
    "Latitude": "Latitude",
    "Longitude": "Longitude",
    "Endereco Completo": "Endereco Completo"

})

df.columns = df.columns.str.strip()

# =========================
# COLUNAS PADRÃO
# =========================

colunas_padrao = {
    "Data": "",
    "Cliente": "Não informado",
    "Responsável Compra": "",
    "WhatsApp": "",
    "Segmento": "Não informado",
    "Cidade": "Não informado",
    "Estado": "MG",
    "Bairro": "",
    "Endereço": "",
    "CEP": "",

    "Concorrente": "Não informado",
    "Tempo Fornecedor": "",
    "Fator Decisão": "",
    "Fornecedor Atende": "",
    "Dificuldade Fornecedor Atual": "",
    "Problemas Atuais": "",

    "Qualidade": 0,
    "Frequência Entrega": 0,
    "Atendimento Atual": 0,

    "Ja Comprou PifPaf": "",
    "PifPaf Preço": "",
    "PifPaf Qualidade": "",
    "PifPaf Entrega": "",
    "PifPaf Atendimento": "",
    "PifPaf Variedade": "",
    "PifPaf Negociação": "",

    "Oportunidade PifPaf": "",
    "Dificuldade PifPaf": "",

    "Categoria Mais Vendida": "Não informado",
    "Volume Mensal": 0,
    "Potencial": "Baixo",
    "Interesse em Trocar": "Não",

    "Status Comercial": "Prospectar",

    # compatibilidade com o código antigo
    "Índice de Atraso": 0,
    "Itens com Ruptura": 0
}

for coluna, valor in colunas_padrao.items():
    if coluna not in df.columns:
        df[coluna] = valor

for coluna, valor in colunas_padrao.items():
    df[coluna] = df[coluna].replace(["", "nan", "None"], valor)

# =========================
# CONVERTER TEXTO EM NÚMERO
# =========================

def converter_volume(valor):
    valor = str(valor).strip()

    if "Até" in valor:
        return 2000
    elif "2.000" in valor and "5.000" in valor:
        return 5000
    elif "5.000" in valor and "10.000" in valor:
        return 10000
    elif "Acima" in valor:
        return 15000
    else:
        valor = (
            valor
            .replace("R$", "")
            .replace(".", "")
            .replace(",", ".")
        )
        return pd.to_numeric(valor, errors="coerce")

df["Volume Mensal"] = df["Volume Mensal"].apply(converter_volume).fillna(0)

colunas_numericas = [
    "Qualidade",
    "Frequência Entrega",
    "Atendimento Atual",
    "Índice de Atraso",
    "Itens com Ruptura"
]

for coluna in colunas_numericas:
    df[coluna] = pd.to_numeric(df[coluna], errors="coerce").fillna(0)

# usar frequência como base de atraso operacional
df["Índice de Atraso"] = 5 - df["Frequência Entrega"]

# usar problemas atuais para estimar ruptura
df["Problemas Atuais"] = df["Problemas Atuais"].fillna("").astype("string").astype(str)

df["Itens com Ruptura"] = df["Problemas Atuais"].apply(
    lambda x: 5 if ("Ruptura" in str(x)) or ("Falta de produtos" in str(x)) else 0
)

# =========================
# COMPARAÇÃO LINHA A LINHA
# =========================

def nota_texto(valor):

    if pd.isna(valor):
        return 0

    valor = str(valor).strip().lower()

    mapa = {
        "péssimo": 1,
        "pessimo": 1,
        "ruim": 2,
        "regular": 3,
        "bom": 4,
        "excelente": 5,
        "ótimo": 5,
        "otimo": 5
    }

    if valor in mapa:
        return mapa[valor]

    try:
        return float(valor)

    except:
        return 0


colunas_avaliacao = [
    "Qualidade",
    "Frequência Entrega",
    "Atendimento Atual",
    "PifPaf Preço",
    "PifPaf Qualidade",
    "PifPaf Entrega",
    "PifPaf Atendimento",
    "PifPaf Variedade",
    "PifPaf Negociação"
]

for coluna in colunas_avaliacao:

    if coluna not in df.columns:
        df[coluna] = 0

    df[coluna] = df[coluna].apply(nota_texto)

    df[coluna] = pd.to_numeric(
        df[coluna],
        errors="coerce"
    ).fillna(0)


df["Score Fornecedor Atual"] = (
    df["Qualidade"] +
    df["Frequência Entrega"] +
    df["Atendimento Atual"]
) / 3


df["Score PifPaf"] = (
    df["PifPaf Preço"] +
    df["PifPaf Qualidade"] +
    df["PifPaf Entrega"] +
    df["PifPaf Atendimento"] +
    df["PifPaf Variedade"] +
    df["PifPaf Negociação"]
) / 6


df["Vantagem PifPaf"] = (
    df["Score PifPaf"] -
    df["Score Fornecedor Atual"]
)


df["Diagnóstico Comparativo"] = df["Vantagem PifPaf"].apply(
    lambda x: "PifPaf melhor posicionada" if x > 0.5
    else "Fornecedor atual melhor posicionado" if x < -0.5
    else "Disputa equilibrada"
)

# =========================
# PADRONIZAR PIFPAF
# =========================

df["Concorrente"] = df["Concorrente"].replace({
    "pifpaf": "PifPaf",
    "Pif paf": "PifPaf",
    "Pif Paf": "PifPaf",
    "pif paf": "PifPaf"
})

# st.success("✅ Dados carregados automaticamente do Google Sheets")

# =========================
# TOPO PREMIUM SEM HTML
# =========================

st.container(border=True)

col_titulo, col_pifpaf, col_status = st.columns([4,1,1])


with col_titulo:

    st.markdown("""
    # EA Comercial Intelligence

    CRM Analítico Comercial e Operacional

    Região: Vespasiano, Confins e São José da Lapa
    """)

st.success("🟢 Status: Online")


   # st.info("Base conectada: Google Forms + Google Sheets")
# =========================
# SIDEBAR / FILTROS
# =========================

st.sidebar.image(
    "assets/logo_ea.png",
    width=180
)

st.sidebar.caption("CRM Analítico Comercial")

st.sidebar.markdown("---")
st.sidebar.subheader("Filtros Estratégicos")

cidades_lista = sorted(df["Cidade"].dropna().unique())
clientes_lista = sorted(df["Cliente"].dropna().unique())
empresas_lista = sorted(df["Concorrente"].dropna().unique())
potenciais_lista = sorted(df["Potencial"].dropna().unique())

cidades = st.sidebar.multiselect(
    "Cidade",
    cidades_lista,
    default=cidades_lista
)

clientes = st.sidebar.multiselect(
    "Cliente",
    clientes_lista,
    default=clientes_lista
)

empresas = st.sidebar.multiselect(
    "Empresa / Concorrente",
    empresas_lista,
    default=empresas_lista
)

potenciais = st.sidebar.multiselect(
    "Potencial",
    potenciais_lista,
    default=potenciais_lista
)

st.sidebar.markdown("---")
st.sidebar.caption("Euler Antonio de Souza")
st.sidebar.caption("Status: Online")

df_filtrado = df[
    (df["Cidade"].isin(cidades)) &
    (df["Cliente"].isin(clientes)) &
    (df["Concorrente"].isin(empresas)) &
    (df["Potencial"].isin(potenciais))
]
# =========================
# ABAS DO SISTEMA
# =========================

aba_geral, aba_clientes, aba_vendas, aba_concorrencia, aba_funil_comercial, aba_estoque, aba_mapa, aba_relatorio, aba_pdf, aba_admin = st.tabs([
    "Visão Geral",
    "CRM Clientes",
    "Vendas Reais",
    "Análise de Concorrência",
    "Funil Comercial",
    "Estoque",
    "Mapa",
    "Relatórios",
    "Importar PDF",
    "Admin"
])

# =========================
# ABA 1 — VISÃO GERAL
# =========================

with aba_geral:

    st.markdown("## 🧠 Central de Inteligência Comercial")
    st.caption("Visão executiva da base de pesquisa, oportunidades, concorrência e prioridade comercial.")

    df_visao = df_filtrado.copy()

    total_clientes = len(df_visao)
    volume_total = df_visao["Volume Mensal"].sum() if "Volume Mensal" in df_visao.columns else 0
    score_fornecedor = df_visao["Score Fornecedor Atual"].mean() if not df_visao.empty else 0
    score_pifpaf = df_visao["Score PifPaf"].mean() if not df_visao.empty else 0
    vantagem_media = df_visao["Vantagem PifPaf"].mean() if not df_visao.empty else 0

    clientes_alto = len(
        df_visao[
            df_visao["Potencial"].astype(str).str.contains("Alto|5", case=False, na=False)
        ]
    ) if not df_visao.empty else 0

    abertura_troca = len(
        df_visao[
            df_visao["Interesse em Trocar"].astype(str).str.contains("Sim|Talvez", case=False, na=False)
        ]
    ) if not df_visao.empty else 0

    # Score comercial de 0 a 100
    df_visao["Score Comercial IA"] = 0

    if not df_visao.empty:
        volume_max = df_visao["Volume Mensal"].max() if df_visao["Volume Mensal"].max() > 0 else 1

        df_visao["Score Comercial IA"] = (
            (df_visao["Volume Mensal"] / volume_max * 25)
            + (df_visao["Score PifPaf"].fillna(0) * 10)
            + ((df_visao["Vantagem PifPaf"].fillna(0) + 5) * 5)
        )

        df_visao.loc[
            df_visao["Interesse em Trocar"].astype(str).str.contains("Sim", case=False, na=False),
            "Score Comercial IA"
        ] += 20

        df_visao.loc[
            df_visao["Interesse em Trocar"].astype(str).str.contains("Talvez", case=False, na=False),
            "Score Comercial IA"
        ] += 10

        df_visao.loc[
            df_visao["Potencial"].astype(str).str.contains("Alto|5", case=False, na=False),
            "Score Comercial IA"
        ] += 15

        df_visao["Score Comercial IA"] = df_visao["Score Comercial IA"].clip(0, 100).round(1)

    score_medio_ia = df_visao["Score Comercial IA"].mean() if not df_visao.empty else 0

    if vantagem_media <= -1:
        risco_concorrencial = "Alto"
    elif vantagem_media < 0:
        risco_concorrencial = "Médio"
    else:
        risco_concorrencial = "Baixo"

    col1, col2, col3, col4, col5 = st.columns(5)

    with col1:
        st.metric("Clientes Mapeados", total_clientes)

    with col2:
        st.metric("Potencial Mensal", f"R$ {volume_total:,.0f}")

    with col3:
        st.metric("Score IA", f"{score_medio_ia:.1f}/100")

    with col4:
        st.metric("Risco Concorrencial", risco_concorrencial)

    with col5:
        st.metric("Alta Prioridade", clientes_alto)

    st.markdown("---")

    col_graf1, col_graf2 = st.columns(2)

    with col_graf1:
        st.subheader("🏆 Top Clientes por Score Comercial")

        if not df_visao.empty:
            ranking_ia = (
                df_visao[[
                    "Cliente",
                    "Cidade",
                    "Concorrente",
                    "Volume Mensal",
                    "Potencial",
                    "Score Comercial IA"
                ]]
                .sort_values(by="Score Comercial IA", ascending=False)
                .head(10)
            )

            fig_score_ia = px.bar(
                ranking_ia,
                x="Cliente",
                y="Score Comercial IA",
                color="Score Comercial IA",
                text="Score Comercial IA",
                title="Prioridade Comercial por Cliente"
            )

            fig_score_ia.update_traces(textposition="outside")
            st.plotly_chart(fig_score_ia, use_container_width=True)
        else:
            st.info("Nenhum dado disponível para ranking.")

    with col_graf2:
        st.subheader("🥧 Participação por Concorrente")

        if not df_visao.empty:
            share_concorrente = (
                df_visao["Concorrente"]
                .value_counts()
                .reset_index()
            )
            share_concorrente.columns = ["Concorrente", "Clientes"]

            fig_share = px.pie(
                share_concorrente,
                names="Concorrente",
                values="Clientes",
                title="Share da Base Mapeada"
            )

            st.plotly_chart(fig_share, use_container_width=True)
        else:
            st.info("Nenhum dado disponível para share.")

    st.markdown("---")

    col_analise1, col_analise2 = st.columns(2)

    with col_analise1:
        st.subheader("📍 Potencial por Cidade")

        if not df_visao.empty:
            potencial_cidade = (
                df_visao.groupby("Cidade")
                .agg(
                    clientes=("Cliente", "count"),
                    volume=("Volume Mensal", "sum"),
                    score_ia=("Score Comercial IA", "mean")
                )
                .reset_index()
                .sort_values(by="volume", ascending=False)
            )

            fig_cidade = px.bar(
                potencial_cidade,
                x="Cidade",
                y="volume",
                color="score_ia",
                text="volume",
                title="Potencial Financeiro por Cidade"
            )

            st.plotly_chart(fig_cidade, use_container_width=True)
        else:
            st.info("Nenhum dado disponível por cidade.")

    with col_analise2:
        st.subheader("⚔️ PifPaf x Fornecedor Atual")

        comparativo = pd.DataFrame({
            "Indicador": ["Fornecedor Atual", "PifPaf"],
            "Score": [score_fornecedor, score_pifpaf]
        })

        fig_comp = px.bar(
            comparativo,
            x="Indicador",
            y="Score",
            color="Indicador",
            text="Score",
            title="Comparativo Médio de Percepção"
        )

        fig_comp.update_traces(texttemplate="%{text:.2f}", textposition="outside")
        fig_comp.update_layout(yaxis=dict(range=[0, 5]))
        st.plotly_chart(fig_comp, use_container_width=True)

    st.markdown("---")

    st.subheader("🧠 Insights Executivos Automáticos")

    if total_clientes == 0:
        st.warning("Nenhum cliente encontrado com os filtros atuais.")
    else:
        if score_pifpaf > score_fornecedor:
            st.success(
                f"A PifPaf está melhor posicionada na média da base filtrada: "
                f"{score_pifpaf:.2f} contra {score_fornecedor:.2f} do fornecedor atual."
            )
        else:
            st.warning(
                f"O fornecedor atual está mais forte na média da base filtrada: "
                f"{score_fornecedor:.2f} contra {score_pifpaf:.2f} da PifPaf."
            )

        if abertura_troca > 0:
            st.info(
                f"Existem {abertura_troca} cliente(s) com abertura para troca ou ampliação de fornecimento."
            )

        if clientes_alto > 0:
            st.success(
                f"Foram identificados {clientes_alto} cliente(s) de alta prioridade comercial."
            )
        else:
            st.info(
                "Nenhum cliente foi classificado como alta prioridade nos filtros atuais."
            )

        if volume_total > 0:
            st.info(
                f"O potencial mensal mapeado nos filtros atuais é de aproximadamente R$ {volume_total:,.0f}."
            )

    st.markdown("---")

    st.subheader("📋 Dados Estratégicos da Pesquisa")

    colunas_tabela = [
        "Data",
        "Cliente",
        "Segmento",
        "Concorrente",
        "Cidade",
        "Volume Mensal",
        "Potencial",
        "Interesse em Trocar",
        "Oportunidade PifPaf",
        "Dificuldade PifPaf",
        "Score Fornecedor Atual",
        "Score PifPaf",
        "Vantagem PifPaf",
        "Score Comercial IA",
        "Diagnóstico Comparativo",
        "Status Comercial"
    ]

    colunas_existentes = [
        coluna for coluna in colunas_tabela
        if coluna in df_visao.columns
    ]

    st.dataframe(
        df_visao[colunas_existentes],
        use_container_width=True,
        hide_index=True
    )

# =========================
# ABA — ANÁLISE DE CONCORRÊNCIA
# =========================

with aba_concorrencia:

    st.markdown(
        '<div class="ea-section-title">Análise de Concorrência</div>',
        unsafe_allow_html=True
    )

    st.markdown("### PifPaf x Fornecedor Atual")

    comparativo_score = pd.DataFrame({
        "Grupo": ["Fornecedor Atual", "PifPaf"],
        "Score Médio": [
            df_filtrado["Score Fornecedor Atual"].mean(),
            df_filtrado["Score PifPaf"].mean()
        ]
    })

    fig_score = px.bar(
        comparativo_score,
        x="Grupo",
        y="Score Médio",
        color="Grupo",
        text="Score Médio",
        title="Score Médio Comparativo"
    )

    fig_score.update_traces(
        texttemplate="%{text:.2f}",
        textposition="outside"
    )

    fig_score.update_layout(
        height=420,
        yaxis=dict(range=[0, 5]),
        template="plotly_dark",
        showlegend=False
    )

    st.plotly_chart(fig_score, use_container_width=True)

    st.markdown("---")
    st.markdown("### Diagnóstico por Cliente")

    tabela_comparativa = df_filtrado[
        [
            "Cliente",
            "Concorrente",
            "Cidade",
            "Score Fornecedor Atual",
            "Score PifPaf",
            "Vantagem PifPaf",
            "Diagnóstico Comparativo"
        ]
    ].sort_values(
        by="Vantagem PifPaf",
        ascending=False
    )

    st.dataframe(
        tabela_comparativa,
        use_container_width=True,
        hide_index=True
    )

    st.markdown("---")
    st.markdown("### Maiores Oportunidades para PifPaf")

    oportunidades = tabela_comparativa[
        tabela_comparativa["Diagnóstico Comparativo"] == "PifPaf melhor posicionada"
    ]

    if oportunidades.empty:
        st.warning("Nenhuma oportunidade clara identificada no momento.")
    else:
        st.success("Clientes onde a PifPaf está melhor posicionada:")
        st.dataframe(
            oportunidades,
            use_container_width=True,
            hide_index=True
        )


# =========================
# ABA — ESTOQUE
# =========================

with aba_estoque:

    st.header("📦 Estoque Inteligente")

    st.info(
        "Envie o arquivo diário de estoque para consultar produtos por código ou descrição. "
        "Depois de salvar, o estoque também fica registrado no banco local."
    )

    arquivo_estoque = st.file_uploader(
        "📤 Envie o arquivo de estoque diário",
        type=["xlsx", "xls"],
        key="upload_estoque_diario"
    )

    if arquivo_estoque:

        try:
            excel_estoque = pd.ExcelFile(arquivo_estoque)

            abas_validas = [
                aba for aba in excel_estoque.sheet_names
                if not str(aba).startswith("_")
            ]

            aba_padrao = "Geral" if "Geral" in abas_validas else abas_validas[0]

            aba_estoque_selecionada = st.selectbox(
                "Selecione a aba do arquivo",
                abas_validas,
                index=abas_validas.index(aba_padrao)
            )

            df_estoque_original = pd.read_excel(
                arquivo_estoque,
                sheet_name=aba_estoque_selecionada
            )

            df_estoque, colunas_estoque = preparar_estoque(df_estoque_original)

            st.success("Estoque carregado com sucesso!")

            col1, col2, col3, col4 = st.columns(4)

            with col1:
                st.metric("Produtos na base", len(df_estoque))

            with col2:
                st.metric(
                    "Estoque Total",
                    f'{df_estoque["_estoque_total"].sum():,.2f}'
                )

            with col3:
                itens_zerados = len(df_estoque[df_estoque["_estoque_total"] <= 0])
                st.metric("Itens zerados", itens_zerados)

            with col4:
                st.metric("Aba lida", aba_estoque_selecionada)

            st.markdown("---")
            st.markdown("---")

            # =========================
            # FILTROS ESTOQUE
            # =========================

            col_filtro1, col_filtro2 = st.columns([2, 1])

            with col_filtro1:

                busca_produto = st.text_input(
                    "🔎 Lupa de pesquisa — digite código, descrição ou categoria do produto"
                )

            with col_filtro2:

                categorias_disponiveis = sorted(
                    df_estoque["_categoria"]
                    .dropna()
                    .astype(str)
                    .unique()
                )

                categoria_selecionada = st.selectbox(
                    "📂 Categoria",
                    ["Todas"] + categorias_disponiveis
                )

            df_resultado_estoque = df_estoque.copy()

            # =========================
            # FILTRO DE PESQUISA
            # =========================

            if busca_produto:

                df_resultado_estoque = df_resultado_estoque[
                    df_resultado_estoque["_material"].str.contains(
                        busca_produto,
                        case=False,
                        na=False
                    )
                    |
                    df_resultado_estoque["_descricao"].str.contains(
                        busca_produto,
                        case=False,
                        na=False
                    )
                    |
                    df_resultado_estoque["_categoria"].str.contains(
                        busca_produto,
                        case=False,
                        na=False
                    )
                ]

            # =========================
            # FILTRO CATEGORIA
            # =========================

            if categoria_selecionada != "Todas":

                df_resultado_estoque = df_resultado_estoque[
                    df_resultado_estoque["_categoria"].astype(str)
                    == categoria_selecionada
                ]

            st.subheader("📋 Resultado da Pesquisa")

            colunas_exibicao_estoque = []
            for coluna in [
                colunas_estoque["centro"],
                colunas_estoque["material"],
                colunas_estoque["descricao"],
                colunas_estoque["categoria"],
                colunas_estoque["total"],
            ]:
                if coluna not in colunas_exibicao_estoque:
                    colunas_exibicao_estoque.append(coluna)

            if not df_resultado_estoque.empty:
                st.dataframe(
                    df_resultado_estoque[colunas_exibicao_estoque],
                    use_container_width=True,
                    hide_index=True
                )

                st.markdown("---")

                if st.button("💾 Salvar estoque diário no banco"):
                    data_importacao = salvar_estoque_no_banco(
                        df_estoque,
                        arquivo_estoque.name,
                        aba_estoque_selecionada
                    )

                    st.success(
                        f"Estoque salvo no banco com sucesso em {data_importacao}."
                    )

            else:
                st.warning("Nenhum produto encontrado para essa pesquisa.")

        except Exception as erro:
            st.error("Não consegui ler o arquivo de estoque.")
            st.exception(erro)

    st.markdown("---")
    st.subheader("🗃️ Último estoque salvo no banco")

    conn = conectar()

    df_estoque_salvo = pd.read_sql_query("""
        SELECT
            data_importacao,
            arquivo_origem,
            aba_origem,
            centro,
            material,
            descricao,
            categoria,
            estoque_total
        FROM estoque_diario
        ORDER BY id DESC
    """, conn)

    conn.close()

    if df_estoque_salvo.empty:
        st.info("Nenhum estoque salvo no banco ainda.")
    else:
        ultima_data_estoque = df_estoque_salvo["data_importacao"].iloc[0]
        df_ultimo_estoque = df_estoque_salvo[
            df_estoque_salvo["data_importacao"] == ultima_data_estoque
        ].copy()

        busca_salva = st.text_input(
            "🔎 Pesquisar no último estoque salvo",
            key="busca_estoque_salvo"
        )

        if busca_salva:
            df_ultimo_estoque = df_ultimo_estoque[
                df_ultimo_estoque["material"].astype(str).str.contains(
                    busca_salva,
                    case=False,
                    na=False
                )
                |
                df_ultimo_estoque["descricao"].astype(str).str.contains(
                    busca_salva,
                    case=False,
                    na=False
                )
                |
                df_ultimo_estoque["categoria"].astype(str).str.contains(
                    busca_salva,
                    case=False,
                    na=False
                )
            ]

        st.caption(f"Última carga: {ultima_data_estoque}")

        st.dataframe(
            df_ultimo_estoque,
            use_container_width=True,
            hide_index=True
        )

# =========================
# ABA — MAPA
# =========================

with aba_mapa:

    st.markdown("## Mapa Comercial")

    df_map = df_filtrado.copy()

    if "Latitude" not in df_map.columns:
        df_map["Latitude"] = None

    if "Longitude" not in df_map.columns:
        df_map["Longitude"] = None

    df_map["Latitude"] = (
        df_map["Latitude"]
        .astype(str)
        .str.replace(",", ".", regex=False)
    )

    df_map["Longitude"] = (
        df_map["Longitude"]
        .astype(str)
        .str.replace(",", ".", regex=False)
    )

    df_map["Latitude"] = pd.to_numeric(
        df_map["Latitude"],
        errors="coerce"
    )

    df_map["Longitude"] = pd.to_numeric(
        df_map["Longitude"],
        errors="coerce"
    )

    df_map = df_map.dropna(
        subset=["Latitude", "Longitude"]
    )

    conn = conectar()

    df_clientes_crm = pd.read_sql_query("""
        SELECT
            codigo_cliente,
            nome,
            cnpj,
            endereco,
            latitude,
            longitude,
            metodo_pagamento
        FROM clientes
        WHERE latitude IS NOT NULL
        AND longitude IS NOT NULL
    """, conn)

    conn.close()

    if df_map.empty and df_clientes_crm.empty:

        st.warning(
            "Nenhuma coordenada encontrada. Importe clientes com endereço válido ou preencha latitude/longitude na base."
        )

    else:

        if not df_map.empty:
            centro_lat = df_map["Latitude"].mean()
            centro_lon = df_map["Longitude"].mean()
        else:
            centro_lat = df_clientes_crm["latitude"].mean()
            centro_lon = df_clientes_crm["longitude"].mean()

        mapa = folium.Map(
            location=[centro_lat, centro_lon],
            zoom_start=11,
            tiles="CartoDB positron"
        )

        # Clientes da pesquisa / Google Sheets
        for _, row in df_map.iterrows():

            atraso = row["Índice de Atraso"]
            ruptura = row["Itens com Ruptura"]
            qualidade = row["Qualidade"]

            risco = "Baixo"
            cor = "green"

            if atraso >= 4 or ruptura >= 5 or qualidade <= 2:
                risco = "Crítico"
                cor = "red"

            elif atraso >= 2 or ruptura >= 3:
                risco = "Atenção"
                cor = "orange"

            folium.CircleMarker(
                location=[
                    row["Latitude"],
                    row["Longitude"]
                ],
                radius=12,
                popup=f"""
                <b>Cliente Pesquisa:</b> {row['Cliente']}<br>
                <b>Cidade:</b> {row['Cidade']}<br>
                <b>Fornecedor:</b> {row['Concorrente']}<br>
                <b>Potencial:</b> {row['Potencial']}<br>
                <b>Risco:</b> {risco}
                """,
                color=cor,
                fill=True,
                fill_color=cor,
                fill_opacity=0.7
            ).add_to(mapa)

        # Clientes reais salvos no CRM via PDF
        for _, row in df_clientes_crm.iterrows():

            folium.Marker(
                location=[
                    row["latitude"],
                    row["longitude"]
                ],
                popup=f"""
                <b>Cliente CRM:</b> {row['nome']}<br>
                <b>Código:</b> {row['codigo_cliente']}<br>
                <b>CNPJ:</b> {row['cnpj']}<br>
                <b>Endereço:</b> {row['endereco']}<br>
                <b>Pagamento:</b> {row['metodo_pagamento']}
                """,
                icon=folium.Icon(
                    color="blue",
                    icon="shopping-cart",
                    prefix="fa"
                )
            ).add_to(mapa)

        st_folium(
            mapa,
            width=1200,
            height=540
        )

# =========================
# ABA 4 — CLIENTES
# =========================

with aba_clientes:

    st.header("👥 CRM de Clientes")

    conn = conectar()

    # =========================
    # CLIENTES CADASTRADOS
    # =========================

    df_clientes = pd.read_sql_query("""
        SELECT *
        FROM clientes
        ORDER BY nome
    """, conn)

    if df_clientes.empty:
        st.warning("Nenhum cliente cadastrado ainda.")

    else:

        cliente_selecionado = st.selectbox(
            "Selecione um cliente",
            df_clientes["nome"].unique()
        )

        dados_cliente = df_clientes[
            df_clientes["nome"] == cliente_selecionado
        ].iloc[0]

        # =========================
        # DADOS PRINCIPAIS
        # =========================

        col1, col2, col3 = st.columns(3)

        with col1:
            st.metric(
                "Código Cliente",
                dados_cliente["codigo_cliente"]
            )

        with col2:
            st.metric(
                "Cidade",
                dados_cliente.get("cidade", "Não informado")
            )

        with col3:
            st.metric(
                "Pagamento",
                dados_cliente["metodo_pagamento"]
            )

        st.markdown("---")

        st.subheader("📋 Dados do Cliente")

        st.write("**Cliente:**", dados_cliente["nome"])
        st.write("**CNPJ:**", dados_cliente["cnpj"])
        st.write("**E-mail:**", dados_cliente["email"])
        st.write("**Endereço:**", dados_cliente["endereco"])

        # =========================
        # PEDIDOS / KPIs AVANÇADOS
        # =========================

        df_pedidos = pd.read_sql_query("""
            SELECT *
            FROM pedidos
            WHERE codigo_cliente = ?
        """, conn, params=(dados_cliente["codigo_cliente"],))

        st.markdown("---")
        st.subheader("📊 Indicadores Avançados do Cliente")

        if not df_pedidos.empty:

            total_pedidos = len(df_pedidos)
            valor_total_cliente = df_pedidos["valor_total"].sum()
            ticket_medio = df_pedidos["valor_total"].mean()
            maior_pedido = df_pedidos["valor_total"].max()

            ultimo_pedido = df_pedidos["numero_pedido"].iloc[-1]

            # =========================
            # PRODUTO FAVORITO
            # =========================

            df_produto_favorito = pd.read_sql_query("""
                SELECT
                    produto,
                    SUM(quantidade) as quantidade_total
                FROM itens_pedido
                WHERE numero_pedido IN (
                    SELECT numero_pedido
                    FROM pedidos
                    WHERE codigo_cliente = ?
                )
                GROUP BY produto
                ORDER BY quantidade_total DESC
                LIMIT 1
            """, conn, params=(dados_cliente["codigo_cliente"],))

            if not df_produto_favorito.empty:
                produto_favorito = df_produto_favorito.iloc[0]["produto"]
            else:
                produto_favorito = "Não identificado"

            # =========================
            # FREQUÊNCIA
            # =========================

            if total_pedidos >= 4:
                frequencia = "Alta"
            elif total_pedidos >= 2:
                frequencia = "Média"
            else:
                frequencia = "Baixa"

            # =========================
            # STATUS DO CLIENTE
            # =========================

            if valor_total_cliente >= 10000:
                status_cliente = "Cliente Estratégico"
            elif valor_total_cliente >= 3000:
                status_cliente = "Cliente Ativo"
            else:
                status_cliente = "Cliente em Desenvolvimento"

            # =========================
            # INTELIGÊNCIA COMERCIAL
            # =========================

            if frequencia == "Alta":
                risco_perda = "Baixo"
            elif frequencia == "Média":
                risco_perda = "Médio"
            else:
                risco_perda = "Atenção"

            if ticket_medio >= 5000:
                potencial_auto = "Alto"
            elif ticket_medio >= 2500:
                potencial_auto = "Médio"
            else:
                potencial_auto = "Baixo"

            if total_pedidos == 1:
                recompra_prevista = "Acompanhar em 7 dias"
            elif frequencia == "Alta":
                recompra_prevista = "Alta chance de recompra"
            else:
                recompra_prevista = "Manter contato ativo"

            if produto_favorito != "Não identificado":
                sugestao_mix = f"Oferecer complemento: {produto_favorito[:18]}"
            else:
                sugestao_mix = "Mapear mix do cliente"

            col1, col2, col3, col4 = st.columns(4)

            with col1:
                st.metric("Total de Pedidos", total_pedidos)

            with col2:
                st.metric("Faturamento Total", f"R$ {valor_total_cliente:,.2f}")

            with col3:
                st.metric("Ticket Médio", f"R$ {ticket_medio:,.2f}")

            with col4:
                st.metric("Maior Pedido", f"R$ {maior_pedido:,.2f}")

            col5, col6, col7, col8 = st.columns(4)

            with col5:
                st.metric("Último Pedido", ultimo_pedido)

            with col6:
                st.metric("Produto Favorito", produto_favorito[:22])

            with col7:
                st.metric("Frequência", frequencia)

            with col8:
                st.metric("Status", status_cliente)

            col9, col10, col11, col12 = st.columns(4)

            with col9:
                st.metric("Risco de Perda", risco_perda)

            with col10:
                st.metric("Potencial Auto", potencial_auto)

            with col11:
                st.metric("Recompra", recompra_prevista)

            with col12:
                st.metric("Sugestão Mix", sugestao_mix[:22])

            st.markdown("---")
            st.subheader("📜 Histórico de Pedidos")

            st.dataframe(
                df_pedidos,
                use_container_width=True,
                hide_index=True
            )

        else:
            st.info("Nenhum pedido encontrado.")
        # =========================
        # ITENS MAIS COMPRADOS
        # =========================

        st.markdown("---")

        st.subheader("📦 Produtos Mais Comprados")

        df_itens = pd.read_sql_query("""
            SELECT
                produto,
                SUM(quantidade) as quantidade_total,
                SUM(total_com_st) as valor_total
            FROM itens_pedido
            WHERE numero_pedido IN (
                SELECT numero_pedido
                FROM pedidos
                WHERE codigo_cliente = ?
            )
            GROUP BY produto
            ORDER BY quantidade_total DESC
        """, conn, params=(dados_cliente["codigo_cliente"],))

        if not df_itens.empty:

            st.dataframe(
                df_itens,
                use_container_width=True,
                hide_index=True
            )

            fig = px.bar(
                df_itens,
                x="produto",
                y="quantidade_total",
                title="Produtos Mais Comprados"
            )

            st.plotly_chart(
                fig,
                use_container_width=True
            )

        else:
            st.info("Nenhum item encontrado.")

    conn.close()
# =========================
# ABA — VENDAS REAIS
# =========================

with aba_vendas:

    st.header("📊 Dashboard Real de Vendas")

    conn = conectar()

    df_vendas = pd.read_sql_query("""
        SELECT
            pedidos.numero_pedido,
            pedidos.codigo_cliente,
            clientes.nome AS cliente,
            clientes.cidade,
            clientes.estado,
            pedidos.valor_total,
            pedidos.origem
        FROM pedidos
        LEFT JOIN clientes
        ON pedidos.codigo_cliente = clientes.codigo_cliente
    """, conn)

    df_produtos = pd.read_sql_query("""
        SELECT
            itens_pedido.numero_pedido,
            pedidos.codigo_cliente,
            clientes.nome AS cliente,
            itens_pedido.codigo_produto,
            itens_pedido.produto,
            itens_pedido.quantidade,
            itens_pedido.preco_unitario,
            itens_pedido.total_com_st
        FROM itens_pedido
        LEFT JOIN pedidos
        ON itens_pedido.numero_pedido = pedidos.numero_pedido
        LEFT JOIN clientes
        ON pedidos.codigo_cliente = clientes.codigo_cliente
    """, conn)

    conn.close()

    if df_vendas.empty:
        st.warning("Nenhum pedido salvo ainda. Importe PDFs para alimentar o dashboard.")

    else:

        faturamento_total = df_vendas["valor_total"].sum()
        total_pedidos = len(df_vendas)
        ticket_medio = df_vendas["valor_total"].mean()
        total_clientes_vendas = df_vendas["codigo_cliente"].nunique()

        col1, col2, col3, col4 = st.columns(4)

        with col1:
            st.metric("Faturamento Total", f"R$ {faturamento_total:,.2f}")

        with col2:
            st.metric("Pedidos", total_pedidos)

        with col3:
            st.metric("Ticket Médio", f"R$ {ticket_medio:,.2f}")

        with col4:
            st.metric("Clientes com Compra", total_clientes_vendas)

        st.markdown("---")

        st.subheader("🏆 Ranking de Clientes por Faturamento")

        ranking_clientes = (
            df_vendas
            .groupby("cliente")["valor_total"]
            .sum()
            .reset_index()
            .sort_values(by="valor_total", ascending=False)
        )

        st.dataframe(
            ranking_clientes,
            use_container_width=True,
            hide_index=True
        )

        fig_clientes = px.bar(
            ranking_clientes,
            x="cliente",
            y="valor_total",
            title="Ranking de Clientes por Faturamento",
            text="valor_total"
        )

        st.plotly_chart(fig_clientes, use_container_width=True)

        st.markdown("---")

        st.subheader("📦 Ranking de Produtos Vendidos")

        if not df_produtos.empty:

            ranking_produtos = (
                df_produtos
                .groupby("produto")
                .agg(
                    quantidade_total=("quantidade", "sum"),
                    valor_total=("total_com_st", "sum")
                )
                .reset_index()
                .sort_values(by="valor_total", ascending=False)
            )

            st.dataframe(
                ranking_produtos,
                use_container_width=True,
                hide_index=True
            )

            fig_produtos = px.bar(
                ranking_produtos,
                x="produto",
                y="valor_total",
                title="Produtos por Faturamento",
                text="valor_total"
            )

            st.plotly_chart(fig_produtos, use_container_width=True)

        else:
            st.info("Nenhum item de pedido salvo ainda.")

        st.markdown("---")

        st.subheader("📋 Pedidos Importados")

        st.dataframe(
            df_vendas,
            use_container_width=True,
            hide_index=True
        )
        # =========================
        # CURVA ABC — PRODUTOS
        # =========================

        st.markdown("---")
        st.subheader("📦 Curva ABC de Produtos")

        if not df_produtos.empty:

            curva_abc_produtos = (
                df_produtos
                .groupby("produto")
                .agg(
                    faturamento=("total_com_st", "sum")
                )
                .reset_index()
                .sort_values(by="faturamento", ascending=False)
            )

            curva_abc_produtos["percentual"] = (
                curva_abc_produtos["faturamento"]
                / curva_abc_produtos["faturamento"].sum()
            ) * 100

            curva_abc_produtos["percentual_acumulado"] = (
                curva_abc_produtos["percentual"].cumsum()
            )

            def classificar_abc(valor):
                if valor <= 80:
                    return "A"
                elif valor <= 95:
                    return "B"
                else:
                    return "C"

            curva_abc_produtos["classe_abc"] = (
                curva_abc_produtos["percentual_acumulado"]
                .apply(classificar_abc)
            )

            st.dataframe(
                curva_abc_produtos,
                use_container_width=True,
                hide_index=True
            )

            fig_abc_produtos = px.bar(
                curva_abc_produtos,
                x="produto",
                y="faturamento",
                color="classe_abc",
                title="Curva ABC de Produtos",
                text="faturamento"
            )

            st.plotly_chart(
                fig_abc_produtos,
                use_container_width=True
            )

        # =========================
        # CURVA ABC — CLIENTES
        # =========================

        st.markdown("---")
        st.subheader("👥 Curva ABC de Clientes")

        curva_abc_clientes = (
            df_vendas
            .groupby("cliente")
            .agg(
                faturamento=("valor_total", "sum")
            )
            .reset_index()
            .sort_values(by="faturamento", ascending=False)
        )

        curva_abc_clientes["percentual"] = (
            curva_abc_clientes["faturamento"]
            / curva_abc_clientes["faturamento"].sum()
        ) * 100

        curva_abc_clientes["percentual_acumulado"] = (
            curva_abc_clientes["percentual"].cumsum()
        )

        curva_abc_clientes["classe_abc"] = (
            curva_abc_clientes["percentual_acumulado"]
            .apply(classificar_abc)
        )

        st.dataframe(
            curva_abc_clientes,
            use_container_width=True,
            hide_index=True
        )

        fig_abc_clientes = px.bar(
            curva_abc_clientes,
            x="cliente",
            y="faturamento",
            color="classe_abc",
            title="Curva ABC de Clientes",
            text="faturamento"
        )

        st.plotly_chart(
            fig_abc_clientes,
            use_container_width=True
        )

    conn.close()

# =========================
# ABA — FUNIL COMERCIAL
# =========================

with aba_funil_comercial:

    st.header("🚀 Funil Comercial")

    conn = conectar()

    df_clientes_pipeline = pd.read_sql_query("""
        SELECT *
        FROM clientes
        ORDER BY nome
    """, conn)

    st.subheader("➕ Adicionar Cliente ao Funil")

    with st.form("pipeline_form"):

        cliente_nome = st.selectbox(
            "Cliente",
            df_clientes_pipeline["nome"].unique()
        )

        dados_pipeline = df_clientes_pipeline[
            df_clientes_pipeline["nome"] == cliente_nome
        ].iloc[0]

        etapa = st.selectbox(
            "Etapa Comercial",
            [
                "Lead",
                "Contato",
                "Negociação",
                "Fechado",
                "Pós-venda"
            ]
        )

        responsavel = st.text_input(
            "Responsável",
            value="Euler Souza"
        )

        observacao = st.text_area(
            "Observações Comerciais"
        )

        salvar_pipeline = st.form_submit_button(
            "Salvar no Funil"
        )

        if salvar_pipeline:

            cursor = conn.cursor()

            cursor.execute("""
                INSERT INTO pipeline (
                    codigo_cliente,
                    cliente,
                    etapa,
                    responsavel,
                    observacao,
                    data_atualizacao
                )
                VALUES (?, ?, ?, ?, ?, ?)
            """, (
                dados_pipeline["codigo_cliente"],
                cliente_nome,
                etapa,
                responsavel,
                observacao,
                datetime.now().strftime("%d/%m/%Y %H:%M")
            ))

            conn.commit()

            st.success("Cliente adicionado ao funil!")

    st.markdown("---")

    st.subheader("📊 Funil Atual")

    df_pipeline = pd.read_sql_query("""
        SELECT *
        FROM pipeline
        ORDER BY id DESC
    """, conn)

    if not df_pipeline.empty:

        col1, col2, col3, col4, col5 = st.columns(5)

        with col1:
            st.metric(
                "Leads",
                len(df_pipeline[df_pipeline["etapa"] == "Lead"])
            )

        with col2:
            st.metric(
                "Contato",
                len(df_pipeline[df_pipeline["etapa"] == "Contato"])
            )

        with col3:
            st.metric(
                "Negociação",
                len(df_pipeline[df_pipeline["etapa"] == "Negociação"])
            )

        with col4:
            st.metric(
                "Fechados",
                len(df_pipeline[df_pipeline["etapa"] == "Fechado"])
            )

        with col5:
            st.metric(
                "Pós-venda",
                len(df_pipeline[df_pipeline["etapa"] == "Pós-venda"])
            )

        st.markdown("---")

        st.dataframe(
            df_pipeline,
            use_container_width=True,
            hide_index=True
        )

        fig_pipeline = px.histogram(
            df_pipeline,
            x="etapa",
            title="Distribuição do Funil",
            color="etapa"
        )

        st.plotly_chart(
            fig_pipeline,
            use_container_width=True
        )

    else:
        st.info("Nenhum cliente no funil ainda.")

    conn.close()

with aba_concorrencia:

    # =========================
    # RANKING DE QUALIDADE
    # =========================

    st.markdown("---")
    st.markdown(
        '<div class="ea-section-title">Ranking de Qualidade</div>',
        unsafe_allow_html=True
    )

    ranking_qualidade = df_filtrado.sort_values(
        by="Qualidade",
        ascending=False
    )

    st.dataframe(
        ranking_qualidade[
            [
                "Cliente",
                "Concorrente",
                "Cidade",
                "Qualidade",
                "Potencial"
            ]
        ],
        use_container_width=True,
        hide_index=True
    )

    # =========================
    # INSIGHTS
    # =========================

    st.markdown("---")
    st.markdown(
        '<div class="ea-section-title">Insights Estratégicos</div>',
        unsafe_allow_html=True
    )

    melhor_qualidade = df_filtrado.loc[df_filtrado["Qualidade"].idxmax()]
    pior_atraso = df_filtrado.loc[df_filtrado["Índice de Atraso"].idxmax()]
    maior_ruptura = df_filtrado.loc[df_filtrado["Itens com Ruptura"].idxmax()]

    st.success(f"✅ Melhor qualidade: {melhor_qualidade['Concorrente']}")
    st.warning(f"⚠️ Maior atraso: {pior_atraso['Concorrente']}")
    st.error(f"🚨 Maior ruptura: {maior_ruptura['Concorrente']}")

    # =========================
    # SCORE COMERCIAL PREMIUM
    # =========================

    st.markdown("---")
    st.markdown("## Ranking de Oportunidades")

    df_score = df_filtrado[
        df_filtrado["Concorrente"] != "PifPaf"
    ].copy()

    if df_score.empty:

        st.warning("Não existem concorrentes suficientes para gerar ranking de oportunidades.")

    else:

        scores = []

        for _, row in df_score.iterrows():

            pontos = 0
            pontos += row["Índice de Atraso"] * 10
            pontos += row["Itens com Ruptura"] * 6

            if row["Qualidade"] <= 2:
                pontos += 20

            if row["Potencial"] == "Alto":
                pontos += 30

            elif row["Potencial"] == "Médio":
                pontos += 15

            if row["Interesse em Trocar"] == "Sim":
                pontos += 35

            scores.append(pontos)

        df_score["Score Comercial"] = scores

        def classificar_prioridade(score):

            if score >= 90:
                return "Alta"

            elif score >= 60:
                return "Média"

            return "Baixa"

        df_score["Prioridade"] = df_score["Score Comercial"].apply(classificar_prioridade)

        ranking_score = df_score.sort_values(
            by="Score Comercial",
            ascending=False
        )

        top_oportunidade = ranking_score.iloc[0]

        st.markdown("### Principal Oportunidade Comercial")

        st.success(f"Cliente: {top_oportunidade['Cliente']}")
        st.write(f"Score Comercial: {top_oportunidade['Score Comercial']}")
        st.write(f"Prioridade: {top_oportunidade['Prioridade']}")
        st.write(f"Concorrente Atual: {top_oportunidade['Concorrente']}")

        st.dataframe(
            ranking_score[
                [
                    "Cliente",
                    "Concorrente",
                    "Cidade",
                    "Potencial",
                    "Score Comercial",
                    "Prioridade"
                ]
            ],
            use_container_width=True,
            hide_index=True
        )

        # =========================
        # COMPARATIVO PIFPAF
        # =========================

        st.markdown("---")

        st.markdown(
            '<div class="ea-section-title">Comparativo PifPaf x Fornecedor Atual</div>',
            unsafe_allow_html=True
        )

        score_fornecedor = df_filtrado["Score Fornecedor Atual"].mean()
        score_pifpaf = df_filtrado["Score PifPaf"].mean()
        vantagem_media = df_filtrado["Vantagem PifPaf"].mean()

        col1, col2, col3 = st.columns(3)

        with col1:
            st.markdown("### Score Fornecedor")
            st.write(f"Fornecedor Atual: {round(score_fornecedor, 2)}")

        with col2:
            st.markdown("### Score PifPaf")
            st.write(f"PifPaf: {round(score_pifpaf, 2)}")

        with col3:
            st.markdown("### Vantagem PifPaf")
            st.write(f"Diferença média: {round(vantagem_media, 2)}")

        if score_pifpaf > score_fornecedor:
            st.success("A PifPaf apresenta melhor posicionamento médio na base analisada.")
        elif score_pifpaf < score_fornecedor:
            st.warning("O fornecedor atual apresenta melhor posicionamento médio que a PifPaf.")
        else:
            st.info("A disputa entre PifPaf e fornecedor atual está equilibrada.")
        # =========================
        # OPORTUNIDADES E AÇÕES RECOMENDADAS
        # =========================

        st.markdown("---")

        st.markdown(
            '<div class="ea-section-title">Oportunidades e Ações Recomendadas</div>',
            unsafe_allow_html=True
        )

        for _, row in df_filtrado.iterrows():

            cliente = row["Cliente"]
            empresa = row["Concorrente"]
            atraso = row["Índice de Atraso"]
            ruptura = row["Itens com Ruptura"]
            qualidade = row["Qualidade"]
            potencial = row["Potencial"]

            categoria = (
                row["Categoria Mais Vendida"]
                if "Categoria Mais Vendida" in df_filtrado.columns
                else "categoria Indicadores Operacionais"
            )

            if empresa == "PifPaf":

                if atraso >= 4 or ruptura >= 5 or qualidade <= 2:

                    st.error(f"""
    Cliente: {cliente}

    Gargalo interno identificado na operação da PifPaf.

    Atraso: {atraso}
    Rupturas: {ruptura}
    Qualidade: {qualidade}

    Ações recomendadas:
    - Revisar rota de entrega
    - Avaliar aumento da frequência de abastecimento
    - Verificar cadeia fria e conservação
    - Reduzir ruptura no cliente
    - Criar plano de acompanhamento comercial prioritário
    """)

                else:

                    st.success(f"""
    Cliente: {cliente}

    A operação da PifPaf apresenta estabilidade operacional.

    Ações recomendadas:
    - Manter nível de serviço
    - Preservar relacionamento comercial
    - Monitorar indicadores preventivamente
    """)

            else:

                if atraso >= 4 or ruptura >= 5 or qualidade <= 2:

                    st.warning(f"""
    Cliente: {cliente}

    Concorrente atual: {empresa}

    Falhas operacionais identificadas no fornecedor atual.

    Atraso: {atraso}
    Rupturas: {ruptura}
    Qualidade: {qualidade}

    Oportunidade para PifPaf:
    - Realizar visita comercial
    - Apresentar proposta PifPaf
    - Destacar estabilidade logística
    - Oferecer mix estratégico de {categoria}
    - Usar falhas do concorrente como argumento comercial
    """)

                elif potencial == "Alto":

                    st.info(f"""
    Cliente: {cliente}

    Concorrente atual: {empresa}

    Cliente de alto potencial, porém com baixa vulnerabilidade operacional no momento.

    Ações recomendadas:
    - Manter relacionamento ativo
    - Monitorar fornecedor atual
    - Mapear oportunidades futuras
    - Acompanhar volume e categoria comprada
    """)

                else:

                    st.success(f"""
    Cliente: {cliente}

    Concorrente atual: {empresa}

    Fornecedor atual apresenta estabilidade operacional.

    Ação recomendada:
    - Monitorar periodicamente
    - Manter cliente no radar comercial
    """)

    # =========================
    # ANÁLISE ESTRATÉGICA AVANÇADA
    # =========================

    st.markdown("---")

    st.markdown(
        '<div class="ea-section-title">Análise Estratégica Avançada</div>',
        unsafe_allow_html=True
    )

    analises = []

    score_fornecedor = df_filtrado["Score Fornecedor Atual"].mean()
    score_pifpaf = df_filtrado["Score PifPaf"].mean()
    vantagem_media = df_filtrado["Vantagem PifPaf"].mean()
    volume_total = df_filtrado["Volume Mensal"].sum()

    clientes_prioritarios = len(
        df_filtrado[df_filtrado["Potencial"] == "Alto"]
    )

    interesse_troca = len(
        df_filtrado[
            df_filtrado["Interesse em Trocar"]
            .astype(str)
            .str.contains("Sim|Talvez", case=False, na=False)
        ]
    )

    pifpaf_melhor = len(
        df_filtrado[
            df_filtrado["Diagnóstico Comparativo"] == "PifPaf melhor posicionada"
        ]
    )

    fornecedor_melhor = len(
        df_filtrado[
            df_filtrado["Diagnóstico Comparativo"] == "Fornecedor atual melhor posicionado"
        ]
    )

    equilibrado = len(
        df_filtrado[
            df_filtrado["Diagnóstico Comparativo"] == "Disputa equilibrada"
        ]
    )

    if score_pifpaf > score_fornecedor:

        analises.append(
            f"A PifPaf apresenta score médio superior ao fornecedor atual "
            f"({round(score_pifpaf, 2)} contra {round(score_fornecedor, 2)}). "
            "Isso indica vantagem competitiva percebida na base analisada."
        )

    elif score_pifpaf < score_fornecedor:

        analises.append(
            f"O fornecedor atual apresenta score médio superior à PifPaf "
            f"({round(score_fornecedor, 2)} contra {round(score_pifpaf, 2)}). "
            "Isso indica necessidade de reforçar proposta comercial, negociação ou percepção de valor."
        )

    else:

        analises.append(
            "A PifPaf e o fornecedor atual apresentam score médio equivalente. "
            "O cenário indica disputa equilibrada."
        )

    analises.append(
        f"Na análise cliente por cliente, a PifPaf está melhor posicionada em "
        f"{pifpaf_melhor} cliente(s), o fornecedor atual lidera em "
        f"{fornecedor_melhor} cliente(s), e há disputa equilibrada em "
        f"{equilibrado} cliente(s)."
    )

    if clientes_prioritarios > 0:

        analises.append(
            f"Existem {clientes_prioritarios} cliente(s) de alto potencial que devem receber atenção comercial prioritária."
        )

    if interesse_troca > 0:

        analises.append(
            f"Foram identificados {interesse_troca} cliente(s) com abertura para troca ou ampliação de fornecimento."
        )

    if vantagem_media > 0.5:

        analises.append(
            "A vantagem média da PifPaf é relevante. A recomendação é acelerar visitas comerciais, apresentar proposta estruturada e explorar os pontos fortes percebidos."
        )

    elif vantagem_media < -0.5:

        analises.append(
            "A desvantagem média da PifPaf indica necessidade de ação corretiva em preço, entrega, atendimento ou negociação antes de buscar conversão direta."
        )

    else:

        analises.append(
            "A vantagem média está próxima do equilíbrio. A recomendação é trabalhar relacionamento, negociação consultiva e acompanhamento contínuo."
        )

    if volume_total >= 10000:

        analises.append(
            "O volume mensal analisado demonstra potencial relevante para crescimento regional."
        )

    for texto in analises:

        st.info(texto)


# =========================
# CONTINUAÇÃO — ANÁLISE DE CONCORRÊNCIA
# =========================

with aba_concorrencia:

    st.markdown("## Inteligência de Concorrência")

    # =========================
    # SCORE DE CONVERSÃO INTELIGENTE
    # =========================

    df_score = df_filtrado.copy()

    df_score["Score Conversão"] = (
        (
            df_score["Vantagem PifPaf"].fillna(0) * 25
        )
        +
        (
            df_score["Score PifPaf"].fillna(0) * 10
        )
    )

    # bônus por interesse

    df_score.loc[
        df_score["Interesse em Trocar"]
        .astype(str)
        .str.contains("Sim", case=False, na=False),

        "Score Conversão"
    ] += 30

    df_score.loc[
        df_score["Interesse em Trocar"]
        .astype(str)
        .str.contains("Talvez", case=False, na=False),

        "Score Conversão"
    ] += 15

    # bônus por potencial

    df_score.loc[
        df_score["Potencial"]
        .astype(str)
        .str.contains("Alto", case=False, na=False),

        "Score Conversão"
    ] += 20

    ranking = df_score.sort_values(
        by="Score Conversão",
        ascending=False
    )

    # =========================
    # FATORES DE DECISÃO
    # =========================

    st.markdown("---")
    st.subheader("Fatores de Decisão de Compra")

    fatores = (
        df_filtrado["Fator Decisão"]
        .astype(str)
        .str.split(",")
        .explode()
        .str.strip()
        .value_counts()
        .reset_index()
    )

    fatores.columns = ["Fator", "Quantidade"]

    fig_fatores = px.bar(
        fatores,
        x="Fator",
        y="Quantidade",
        color="Quantidade",
        template="plotly_dark",
        title="O que mais pesa na decisão de compra"
    )

    st.plotly_chart(fig_fatores, use_container_width=True)

    # =========================
    # PROBLEMAS DOS CONCORRENTES
    # =========================

    st.markdown("---")
    st.subheader("Principais Problemas dos Concorrentes")

    problemas = (
        df_filtrado["Problemas Atuais"]
        .astype(str)
        .str.split(",")
        .explode()
        .str.strip()
        .value_counts()
        .reset_index()
    )

    problemas.columns = ["Problema", "Quantidade"]

    fig_problemas = px.bar(
        problemas,
        x="Problema",
        y="Quantidade",
        color="Quantidade",
        template="plotly_dark",
        title="Problemas identificados nos fornecedores"
    )

    st.plotly_chart(fig_problemas, use_container_width=True)

    # =========================
    # OPORTUNIDADES PIFPAF
    # =========================

    st.markdown("---")
    st.subheader("O que faria comprar mais da PifPaf")

    oportunidades = (
        df_filtrado["Oportunidade PifPaf"]
        .astype(str)
        .str.split(",")
        .explode()
        .str.strip()
        .value_counts()
        .reset_index()
    )

    oportunidades.columns = ["Oportunidade", "Quantidade"]

    fig_oportunidades = px.bar(
        oportunidades,
        x="Oportunidade",
        y="Quantidade",
        color="Quantidade",
        template="plotly_dark",
        title="Oportunidades de crescimento PifPaf"
    )

    st.plotly_chart(fig_oportunidades, use_container_width=True)

    # =========================
    # DIFICULDADES PIFPAF
    # =========================

    st.markdown("---")
    st.subheader("Principais Dificuldades com a PifPaf")

    dificuldades = (
        df_filtrado["Dificuldade PifPaf"]
        .astype(str)
        .str.split(",")
        .explode()
        .str.strip()
        .value_counts()
        .reset_index()
    )

    dificuldades.columns = ["Dificuldade", "Quantidade"]

    fig_dificuldades = px.bar(
        dificuldades,
        x="Dificuldade",
        y="Quantidade",
        color="Quantidade",
        template="plotly_dark",
        title="Principais dificuldades com a PifPaf"
    )

    st.plotly_chart(fig_dificuldades, use_container_width=True)

    # =========================
    # PARTICIPAÇÃO DOS CONCORRENTES
    # =========================

    st.markdown("---")
    st.subheader("Participação dos Concorrentes")

    ranking_concorrentes = (
        df_filtrado["Concorrente"]
        .value_counts()
        .reset_index()
    )

    ranking_concorrentes.columns = ["Concorrente", "Clientes"]

    fig_concorrentes = px.pie(
        ranking_concorrentes,
        names="Concorrente",
        values="Clientes",
        template="plotly_dark",
        title="Participação dos fornecedores"
    )

    st.plotly_chart(fig_concorrentes, use_container_width=True)

    # =========================
    # POTENCIAL POR CIDADE
    # =========================

    st.markdown("---")
    st.subheader("Potencial Comercial por Cidade")

    potencial_cidade = (
        df_filtrado.groupby("Cidade")["Cliente"]
        .count()
        .reset_index()
    )

    potencial_cidade.columns = ["Cidade", "Clientes"]

    fig_cidade = px.bar(
        potencial_cidade,
        x="Cidade",
        y="Clientes",
        color="Clientes",
        template="plotly_dark",
        title="Clientes mapeados por cidade"
    )

    st.plotly_chart(fig_cidade, use_container_width=True)

    # =========================
    # RADAR — PERCEPÇÃO PIFPAF
    # =========================

    st.markdown("---")
    st.subheader("Radar de Percepção da PifPaf")

    colunas_radar = [
        "PifPaf Preço",
        "PifPaf Qualidade",
        "PifPaf Entrega",
        "PifPaf Atendimento",
        "PifPaf Variedade",
        "PifPaf Negociação"
    ]

    colunas_existentes_radar = [
        coluna for coluna in colunas_radar
        if coluna in df_filtrado.columns
    ]

    if len(colunas_existentes_radar) > 0:

        df_radar = df_filtrado[colunas_existentes_radar].copy()

        for coluna in colunas_existentes_radar:
            df_radar[coluna] = pd.to_numeric(
                df_radar[coluna],
                errors="coerce"
            ).fillna(0)

        radar_media = df_radar.mean().reset_index()
        radar_media.columns = ["Indicador", "Nota"]

        fig_radar = px.line_polar(
            radar_media,
            r="Nota",
            theta="Indicador",
            line_close=True,
            template="plotly_dark",
            title="Percepção média da PifPaf"
        )

        fig_radar.update_traces(fill="toself")

        st.plotly_chart(fig_radar, use_container_width=True)

    # =========================
    # HEATMAP
    # =========================

    st.markdown("---")
    st.subheader("Mapa de Força Comercial por Cidade e Concorrente")

    heatmap = (
        df_filtrado
        .groupby(["Cidade", "Concorrente"])["Cliente"]
        .count()
        .reset_index()
    )

    heatmap.columns = ["Cidade", "Concorrente", "Clientes"]

    fig_heatmap = px.density_heatmap(
        heatmap,
        x="Concorrente",
        y="Cidade",
        z="Clientes",
        template="plotly_dark",
        title="Concentração de clientes por cidade e fornecedor"
    )

    st.plotly_chart(fig_heatmap, use_container_width=True)

    # =========================
    # FUNIL COMERCIAL
    # =========================

    st.markdown("---")
    st.subheader("Funil Comercial")

    total_mapeados = len(df_filtrado)

    interessados = len(
        df_filtrado[
            df_filtrado["Interesse em Trocar"]
            .astype(str)
            .str.contains("Sim|Talvez", case=False, na=False)
        ]
    )

    alto_potencial = len(
        df_filtrado[
            df_filtrado["Potencial"]
            .astype(str)
            .str.contains("Alto", case=False, na=False)
        ]
    )

    oportunidades_qtd = len(
        df_filtrado[
            df_filtrado["Oportunidade PifPaf"].astype(str).str.len() > 3
        ]
    )

    funil = pd.DataFrame({
        "Etapa": [
            "Clientes mapeados",
            "Com abertura",
            "Alto potencial",
            "Oportunidade identificada"
        ],
        "Quantidade": [
            total_mapeados,
            interessados,
            alto_potencial,
            oportunidades_qtd
        ]
    })

    fig_funil = px.funnel(
        funil,
        x="Quantidade",
        y="Etapa",
        template="plotly_dark",
        title="Funil de oportunidades comerciais"
    )

    st.plotly_chart(fig_funil, use_container_width=True)

    # =========================
    # MATRIZ DE OPORTUNIDADE
    # =========================

    st.markdown("---")
    st.subheader("Matriz de Oportunidade Comercial")

    st.dataframe(
        ranking[
            [
                "Cliente",
                "Cidade",
                "Concorrente",
                "Volume Mensal",
                "Score Conversão",
                "Diagnóstico Comparativo"
            ]
        ],
        use_container_width=True,
        hide_index=True
    )

    # =========================
    # RANKING DE CONVERSÃO
    # =========================

    st.markdown("---")
    st.subheader("Ranking de Conversão Comercial")

    st.dataframe(
        ranking[
            [
                "Cliente",
                "Cidade",
                "Concorrente",
                "Potencial",
                "Interesse em Trocar",
                "Score Conversão"
            ]
        ],
        use_container_width=True,
        hide_index=True
    )
# =========================
# ABA 5 — RELATÓRIO / APRESENTAÇÃO
# =========================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from io import BytesIO

with aba_relatorio:

    st.markdown(
        '<div class="ea-section-title">Apresentação Executiva</div>',
        unsafe_allow_html=True
    )

    if st.button("Gerar Apresentação PowerPoint"):

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        vermelho = RGBColor(153, 27, 27)
        preto = RGBColor(17, 24, 39)
        cinza = RGBColor(107, 114, 128)

        def add_title(slide, title, subtitle=""):
            box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
            tf = box.text_frame
            tf.text = title
            p = tf.paragraphs[0]
            p.font.size = Pt(30)
            p.font.bold = True
            p.font.color.rgb = preto

            if subtitle:
                sub = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(11.5), Inches(0.4))
                sub_tf = sub.text_frame
                sub_tf.text = subtitle
                sub_tf.paragraphs[0].font.size = Pt(14)
                sub_tf.paragraphs[0].font.color.rgb = cinza

        def add_metric(slide, x, y, title, value):
            shape = slide.shapes.add_shape(
                1, Inches(x), Inches(y), Inches(2.4), Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
            shape.line.color.rgb = RGBColor(220, 220, 220)

            tf = shape.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(11)
            tf.paragraphs[0].font.color.rgb = cinza

            p = tf.add_paragraph()
            p.text = str(value)
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.color.rgb = vermelho

        total_clientes = len(df_filtrado)
        volume_total = df_filtrado["Volume Mensal"].sum()
        score_fornecedor = round(df_filtrado["Score Fornecedor Atual"].mean(), 2)
        score_pifpaf = round(df_filtrado["Score PifPaf"].mean(), 2)
        vantagem_media = round(df_filtrado["Vantagem PifPaf"].mean(), 2)

        oportunidades = df_filtrado[
            df_filtrado["Diagnóstico Comparativo"] == "PifPaf melhor posicionada"
        ]

        gargalos = df_filtrado[
            df_filtrado["Diagnóstico Comparativo"] == "Fornecedor atual melhor posicionado"
        ]

        equilibrados = df_filtrado[
            df_filtrado["Diagnóstico Comparativo"] == "Disputa equilibrada"
        ]

        # SLIDE 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "EA Comercial Intelligence", "Apresentação Executiva — Análise PifPaf x Fornecedor Atual")

        tx = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.0))
        tf = tx.text_frame
        tf.text = "Diagnóstico comercial, competitivo e operacional com base nas respostas do Forms."
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = preto

        add_metric(slide, 0.8, 4.2, "Clientes", total_clientes)
        add_metric(slide, 3.4, 4.2, "Score Fornecedor", score_fornecedor)
        add_metric(slide, 6.0, 4.2, "Score PifPaf", score_pifpaf)
        add_metric(slide, 8.6, 4.2, "Vantagem PifPaf", vantagem_media)

        # SLIDE 2
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Resumo Estratégico", "Leitura geral da base analisada")

        resumo = (
            f"Foram analisados {total_clientes} clientes.\n\n"
            f"Score médio do fornecedor atual: {score_fornecedor}\n"
            f"Score médio da PifPaf: {score_pifpaf}\n"
            f"Vantagem média da PifPaf: {vantagem_media}\n\n"
            f"Clientes onde a PifPaf está melhor posicionada: {len(oportunidades)}\n"
            f"Clientes onde o fornecedor atual lidera: {len(gargalos)}\n"
            f"Clientes em disputa equilibrada: {len(equilibrados)}\n\n"
            f"Volume mensal estimado: R$ {volume_total:,.0f}"
        )

        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.8))
        tf = box.text_frame
        tf.text = resumo
        for p in tf.paragraphs:
            p.font.size = Pt(18)
            p.font.color.rgb = preto

        # SLIDE 3
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Comparativo Geral", "Fornecedor Atual x PifPaf")

        rows = 3
        cols = 3
        table = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.8), Inches(11), Inches(1.8)).table

        headers = ["Indicador", "Fornecedor Atual", "PifPaf"]
        values = [
            ["Score Médio", score_fornecedor, score_pifpaf],
            ["Diferença", "-", vantagem_media]
        ]

        for i, h in enumerate(headers):
            table.cell(0, i).text = h

        for r, row in enumerate(values, start=1):
            for c, val in enumerate(row):
                table.cell(r, c).text = str(val)

        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(14)

        # SLIDE 4
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Ranking Comparativo por Cliente", "Clientes ordenados pela vantagem da PifPaf")

        ranking = df_filtrado[
            [
                "Cliente",
                "Concorrente",
                "Cidade",
                "Score Fornecedor Atual",
                "Score PifPaf",
                "Vantagem PifPaf",
                "Diagnóstico Comparativo"
            ]
        ].sort_values(by="Vantagem PifPaf", ascending=False).head(8)

        table = slide.shapes.add_table(
            len(ranking) + 1,
            7,
            Inches(0.3),
            Inches(1.5),
            Inches(12.7),
            Inches(4.8)
        ).table

        headers = list(ranking.columns)

        for c, h in enumerate(headers):
            table.cell(0, c).text = h

        for r, (_, row) in enumerate(ranking.iterrows(), start=1):
            for c, h in enumerate(headers):
                table.cell(r, c).text = str(round(row[h], 2)) if isinstance(row[h], float) else str(row[h])

        for row in table.rows:
            for cell in row.cells:
                cell.text_frame.paragraphs[0].font.size = Pt(8)

        # SLIDE 5
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Maiores Oportunidades", "Clientes onde a PifPaf está melhor posicionada")

        texto = ""
        if oportunidades.empty:
            texto = "Nenhuma oportunidade clara identificada no momento."
        else:
            for _, row in oportunidades.head(6).iterrows():
                texto += f"• {row['Cliente']} — {row['Cidade']} | Vantagem: {round(row['Vantagem PifPaf'], 2)}\n"

        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.8))
        tf = box.text_frame
        tf.text = texto
        for p in tf.paragraphs:
            p.font.size = Pt(18)

        # SLIDE 6
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Gargalos Estratégicos", "Clientes onde o fornecedor atual ainda lidera")

        texto = ""
        if gargalos.empty:
            texto = "Nenhum gargalo estratégico crítico identificado."
        else:
            for _, row in gargalos.head(6).iterrows():
                texto += f"• {row['Cliente']} — {row['Concorrente']} | Diferença: {round(row['Vantagem PifPaf'], 2)}\n"

        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.8), Inches(4.8))
        tf = box.text_frame
        tf.text = texto
        for p in tf.paragraphs:
            p.font.size = Pt(18)

        # SLIDE 7
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Plano de Ação Comercial", "Recomendações práticas")

        plano = """
1. Priorizar clientes onde a PifPaf já aparece melhor posicionada.
2. Criar abordagem consultiva para clientes em disputa equilibrada.
3. Para clientes onde o fornecedor atual lidera, reforçar preço, entrega, variedade e negociação.
4. Trabalhar propostas por categoria de maior volume.
5. Acompanhar evolução dos scores a cada nova pesquisa.
"""

        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.8), Inches(4.8))
        tf = box.text_frame
        tf.text = plano
        for p in tf.paragraphs:
            p.font.size = Pt(19)

        # SLIDE 8
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, "Conclusão Estratégica", "EA Comercial Intelligence")

        conclusao = """
A análise permite identificar com clareza onde a PifPaf possui vantagem competitiva,
onde há disputa equilibrada e onde o fornecedor atual ainda está mais forte.

O modelo transforma respostas do Forms em inteligência comercial prática,
orientando visitas, negociações e priorização de clientes.
"""

        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(11.8), Inches(4.8))
        tf = box.text_frame
        tf.text = conclusao
        for p in tf.paragraphs:
            p.font.size = Pt(22)
            p.font.color.rgb = preto

        arquivo_pptx = BytesIO()
        prs.save(arquivo_pptx)
        arquivo_pptx.seek(0)

        st.download_button(
            label="Baixar Apresentação PowerPoint",
            data=arquivo_pptx,
            file_name="apresentacao_ea_intelligence.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

# =========================
# IMPORTAR PDF
# =========================

with aba_pdf:

    st.header("📄 Importar Pedido por PDF")

    arquivo = st.file_uploader(
        "Envie o PDF do pedido",
        type=["pdf"]
    )

    if arquivo:
        dados = extrair_dados_pedido(arquivo)

        st.success("PDF lido com sucesso!")

        st.subheader("Dados identificados")

        st.write("**Pedido SAP:**", dados["numero_pedido"])
        st.write("**Código Cliente:**", dados["codigo_cliente"])
        st.write("**Cliente:**", dados["nome"])
        st.write("**E-mail:**", dados["email"])
        st.write("**CNPJ:**", dados["cnpj"])
        st.write("**Endereço:**", dados["endereco"])
        st.write("**Pagamento:**", dados["metodo_pagamento"])
        st.write("**Valor Total:**", f'R$ {dados["valor_total"]:,.2f}')

        st.subheader("📦 Itens do Pedido")

        if dados["itens"]:
            df_itens_pdf = pd.DataFrame(dados["itens"])

            st.dataframe(
                df_itens_pdf,
                use_container_width=True,
                hide_index=True
            )
        else:
            st.warning("Nenhum item foi identificado no PDF.")

        if st.button("Salvar no CRM"):
            conn = conectar()
            cursor = conn.cursor()

            latitude, longitude = geocodificar_endereco(
                dados["endereco"]
            )

            cursor.execute("""
                INSERT OR REPLACE INTO clientes (
                    codigo_cliente,
                    nome,
                    email,
                    cnpj,
                    endereco,
                    latitude,
                    longitude,
                    metodo_pagamento
                ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """, (
                dados["codigo_cliente"],
                dados["nome"],
                dados["email"],
                dados["cnpj"],
                dados["endereco"],
                latitude,
                longitude,
                dados["metodo_pagamento"]
            ))

            cursor.execute("""
                INSERT OR REPLACE INTO pedidos (
                    numero_pedido,
                    codigo_cliente,
                    valor_total,
                    origem
                ) VALUES (?, ?, ?, ?)
            """, (
                dados["numero_pedido"],
                dados["codigo_cliente"],
                dados["valor_total"],
                "PDF"
            ))

            # Remove itens antigos desse pedido para evitar duplicidade
            cursor.execute(
                "DELETE FROM itens_pedido WHERE numero_pedido = ?",
                (dados["numero_pedido"],)
            )

            # Salva os itens do pedido
            for item in dados["itens"]:
                cursor.execute("""
                    INSERT INTO itens_pedido (
                        numero_pedido,
                        codigo_produto,
                        produto,
                        quantidade,
                        preco_unitario,
                        total_sem_st,
                        st,
                        total_com_st
                    ) VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """, (
                    dados["numero_pedido"],
                    item["codigo_produto"],
                    item["produto"],
                    item["quantidade"],
                    item["preco_unitario"],
                    item["total_sem_st"],
                    item["st"],
                    item["total_com_st"]
                ))

            conn.commit()
            conn.close()

            if latitude and longitude:
                st.success("Cliente, pedido, itens e localização salvos com sucesso!")
            else:
                st.warning("Cliente, pedido e itens salvos. A localização não foi encontrada automaticamente.")

# =========================
# ABA — ADMIN / MANUTENÇÃO
# =========================

with aba_admin:

    st.header("⚙️ Administração do Sistema")
    st.warning("Área de manutenção. Use com cuidado.")

    conn = conectar()
    cursor = conn.cursor()

    st.subheader("🧾 Gerenciar Pedidos")

    df_pedidos_admin = pd.read_sql_query("""
        SELECT 
            pedidos.numero_pedido,
            pedidos.codigo_cliente,
            clientes.nome AS cliente,
            pedidos.valor_total,
            pedidos.origem
        FROM pedidos
        LEFT JOIN clientes
        ON pedidos.codigo_cliente = clientes.codigo_cliente
        ORDER BY pedidos.numero_pedido DESC
    """, conn)

    if df_pedidos_admin.empty:
        st.info("Nenhum pedido cadastrado.")
    else:
        st.dataframe(df_pedidos_admin, use_container_width=True, hide_index=True)

        pedido_excluir = st.selectbox(
            "Selecione o pedido para excluir",
            df_pedidos_admin["numero_pedido"].astype(str).tolist()
        )

        confirmar_pedido = st.checkbox("Confirmo que desejo excluir este pedido e seus itens")

        if st.button("🗑️ Excluir pedido selecionado"):
            if confirmar_pedido:
                cursor.execute(
                    "DELETE FROM itens_pedido WHERE numero_pedido = ?",
                    (pedido_excluir,)
                )

                cursor.execute(
                    "DELETE FROM pedidos WHERE numero_pedido = ?",
                    (pedido_excluir,)
                )

                conn.commit()
                st.success(f"Pedido {pedido_excluir} excluído com sucesso.")
                st.rerun()
            else:
                st.error("Marque a confirmação antes de excluir.")

    st.markdown("---")

    st.subheader("📦 Gerenciar Estoque Diário")

    df_cargas_estoque = pd.read_sql_query("""
        SELECT
            data_importacao,
            arquivo_origem,
            aba_origem,
            COUNT(*) AS itens
        FROM estoque_diario
        GROUP BY data_importacao, arquivo_origem, aba_origem
        ORDER BY data_importacao DESC
    """, conn)

    if df_cargas_estoque.empty:
        st.info("Nenhuma carga de estoque salva.")
    else:
        st.dataframe(df_cargas_estoque, use_container_width=True, hide_index=True)

        carga_excluir = st.selectbox(
            "Selecione a carga de estoque para excluir",
            df_cargas_estoque["data_importacao"].astype(str).tolist()
        )

        confirmar_estoque = st.checkbox("Confirmo que desejo excluir esta carga de estoque")

        if st.button("🗑️ Excluir carga selecionada"):
            if confirmar_estoque:
                cursor.execute(
                    "DELETE FROM estoque_diario WHERE data_importacao = ?",
                    (carga_excluir,)
                )

                conn.commit()
                st.success(f"Carga de estoque {carga_excluir} excluída com sucesso.")
                st.rerun()
            else:
                st.error("Marque a confirmação antes de excluir.")

        st.markdown("---")

        confirmar_limpeza = st.checkbox("Confirmo que desejo manter somente a última carga de estoque")

        if st.button("🧹 Limpar estoques antigos e manter apenas o último"):
            if confirmar_limpeza:
                ultima_carga = df_cargas_estoque["data_importacao"].iloc[0]

                cursor.execute(
                    "DELETE FROM estoque_diario WHERE data_importacao != ?",
                    (ultima_carga,)
                )

                conn.commit()
                st.success("Estoques antigos removidos. Última carga mantida.")
                st.rerun()
            else:
                st.error("Marque a confirmação antes de limpar.")

    conn.close()
